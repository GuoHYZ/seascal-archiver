from pathlib import Path
from pptx import Presentation
import traceback


def extract_shape_texts(shape, parent_top=0, parent_left=0):
    """
    递归提取形状中的所有文本，返回 [(top, left, text), ...] 列表。
    支持文本框、表格、组合形状等。
    """
    results = []

    try:
        # ---- 处理组合形状 (GroupShape) ----
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for child in shape.shapes:
                results.extend(
                    extract_shape_texts(
                        child,
                        parent_top + (shape.top if hasattr(shape, "top") else 0),
                        parent_left + (shape.left if hasattr(shape, "left") else 0)
                    )
                )
            return results

        # 确定位置（top, left），某些形状可能没有位置属性
        top = parent_top
        left = parent_left
        try:
            top += shape.top or 0
        except AttributeError:
            pass
        try:
            left += shape.left or 0
        except AttributeError:
            pass

        # ---- 处理表格 ----
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        cell_top = top
                        cell_left = left
                        try:
                            cell_top += cell.top or 0
                        except AttributeError:
                            pass
                        try:
                            cell_left += cell.left or 0
                        except AttributeError:
                            pass
                        results.append((cell_top, cell_left, cell_text))
            return results

        # ---- 处理文本框/含文本的形状 ----
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:
                results.append((top, left, text))

    except Exception:
        # 单个形状解析失败不影响其他形状
        pass

    return results


def extract_slide_text(slide):
    """
    提取单张幻灯片中所有文本，按位置（上->下，左->右）排序返回纯文本列表。
    """
    texts = []

    for shape in slide.shapes:
        texts.extend(extract_shape_texts(shape))

    # 根据页面位置排序：先按 top，再按 left
    texts.sort(key=lambda x: (x[0], x[1]))

    return [text for _, _, text in texts]


def ppt_to_txt(ppt_file, txt_file):
    """
    将单个 PPTX 文件转换为 TXT 文件。
    ppt_file: str 或 pathlib.Path
    txt_file: str 或 pathlib.Path
    """
    ppt_file = Path(ppt_file)
    txt_file = Path(txt_file)

    # 自动创建输出目录（包括子目录）
    txt_file.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(ppt_file))

    with open(txt_file, "w", encoding="utf-8-sig", errors="replace") as f:
        f.write(f"文件：{ppt_file.name}\n")
        f.write("=" * 60 + "\n")

        for index, slide in enumerate(prs.slides, start=1):
            f.write(f"\n========== 第 {index} 页 ==========\n\n")

            try:
                texts = extract_slide_text(slide)
            except Exception as e:
                f.write(f"[错误] 无法提取本页内容: {e}\n\n")
                continue

            if not texts:
                f.write("(本页无可提取文本)\n\n")

            for text in texts:
                f.write(text + "\n\n")


def batch_convert(input_dir, output_dir):
    """
    批量将 input_dir 下所有 PPTX 文件转换为 TXT，输出到 output_dir。
    """
    input_dir = Path(input_dir)
    output_dir = Path(output_dir)

    # 检查输入目录是否存在
    if not input_dir.is_dir():
        print(f"[错误] 输入目录不存在: {input_dir}")
        return

    output_dir.mkdir(exist_ok=True)

    ppt_files = list(input_dir.glob("*.pptx"))

    if not ppt_files:
        print(f"[警告] 目录 {input_dir} 中未找到任何 PPTX 文件")
        return

    success = 0
    failed = 0
    total = len(ppt_files)

    for i, ppt in enumerate(ppt_files, start=1):
        # 跳过 office 临时文件
        if ppt.name.startswith("~$"):
            print(f"[{i}/{total}] [跳过] {ppt.name} (临时文件)")
            continue

        txt = output_dir / (ppt.stem + ".txt")

        try:
            ppt_to_txt(ppt, txt)
            success += 1
            print(f"[{i}/{total}] [OK] {ppt.name}")
        except Exception as e:
            failed += 1
            print(f"[{i}/{total}] [FAIL] {ppt.name}")
            print(f"       原因: {e}")
            traceback.print_exc()

    print(f"\n处理结束 - 共 {total} 个文件")
    print(f"成功: {success}, 失败: {failed}")


def convert(source_path, dest_path):
    """
    统一转换接口。
    source_path: str 或 Path — 源 PPTX 文件路径
    dest_path:   str 或 Path — 目标 TXT 文件路径
    返回: int — 提取的总字符数
    """
    ppt_to_txt(source_path, dest_path)

    # 返回提取的字符数
    dest_path = Path(dest_path)
    if dest_path.exists():
        return len(dest_path.read_text(encoding="utf-8-sig", errors="replace"))
    return 0


if __name__ == "__main__":
    batch_convert("./PPT", "./TXT")
